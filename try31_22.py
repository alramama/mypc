import pandas as pd
import re
extracted_file =[]
extracted_file1 =[]
extracted_file2 =[]

def counting1():

    file = pd.read_excel('Jubail Desalination Plant1.xlsx')
    df_sheet_index = pd.read_excel('Jubail Desalination Plant.xlsx', sheet_name=0)
    #print(df_sheet_index)
    data = pd.DataFrame(df_sheet_index)
    rows = data.to_numpy()
    #print(rows)
    print("row",len(rows))
    for i in rows:
        new2 = [i[0],i[2]]
        #extracted_file.append(i[0])
        #print(i)
        if type(i[0]) == str:
            result = (re.split("[A][0-9].[0-9].[0-9]", i[0]))
            result10 = (re.findall("[A][0-9].[0-9].[0-9]", i[0]))
            #print(result10)
            #extracted_file.append(result)
            if len(result) == 2:
                result1 = (re.split("^[0-9]",result[1]))
                #extracted_file.append(result1)
                #print(result1)
                if len(result1) == 1:
                    #print(result1[0])
                    #print(len('Door Monitoring Contact (DMC)'))
                    new2 = (result1[0].strip(),i[1],i[2],result10)
                    new0 = (result1[0].strip())
                    extracted_file.append(new0)
                    extracted_file1.append(new2)
                elif len(result1) == 2:
                    new2 = result1[1].strip(),i[1],i[2],result10
                    new1 = (result1[1].strip())
                    extracted_file.append(new1)
                    extracted_file1.append(new2)

                    #print(len('Door Monitoring Contact (DMC)'))

    #print(extracted_file.sort(reverse=True))
    from collections import Counter
    MyList = ["a", "b", "a", "c", "c", "a", "c"]
    print(len(extracted_file))
    c = Counter(extracted_file)
    #print(len(extracted_file))
    c1 = dict(c)
    string = "Card Reader (CRPK)"
    for key in  c1.keys():
        for value in extracted_file1:
            if value[0] == key:
                extracted_file2.append(((value)))
                #print(value[2])
    data2 = pd.DataFrame(extracted_file2)
    #data2.to_excel('newlist7.xlsx')

    print(extracted_file2)
    import collections
    import matplotlib.pyplot as plt
    l = ['a', 'b', 'b', 'b', 'c']
    # w = collections.Counter(l)
    #w = c
    #plt.bar(w.keys(), w.values())
    #plt.show()

#counting1()

def counting():

    file = pd.read_excel('Jubail Desalination Plant1.xlsx')
    df_sheet_index = pd.read_excel('Jubail Desalination Plant.xlsx', sheet_name=0)
    #print(df_sheet_index)
    data = pd.DataFrame(df_sheet_index)
    rows = data.to_numpy()
    #print(rows)
    print("row",len(rows))
    for i in rows:
        new2 = [i[0],i[2]]
        #extracted_file.append(i[0])
        #print(i)
        if type(i[0]) == str:
            result = re.split("[A][0-9].[0-9].[0-9]", i[0])
            result10 = re.findall("[A][0-9].[0-9].[0-9]", i[0])
            #print(result10)
            #extracted_file.append(result)
            if len(result) == 2:
                result1 = re.split("^[0-9]",result[1])
                #extracted_file.append(result1)
                #print(result1)
                if len(result1) == 1:
                    #print(result1[0])
                    #print(len('Door Monitoring Contact (DMC)'))
                    new2 = result1[0].strip(),i[1],i[2],result10
                    new0 = result1[0].strip()
                    extracted_file.append(new0)
                    extracted_file1.append(new2)
                elif len(result1) == 2:
                    new2 = result1[1].strip(),i[1],i[2],result10
                    new1 =result1[1].strip()
                    extracted_file.append(new1)
                    extracted_file1.append(new2)

                    #print(len('Door Monitoring Contact (DMC)'))

    #print(extracted_file.sort(reverse=True))
    from collections import Counter
    MyList = ["a", "b", "a", "c", "c", "a", "c"]
    print(len(extracted_file))
    c = Counter(extracted_file)
    #print(len(extracted_file))
    c1 = dict(c)
    string = "Card Reader (CRPK)"
    count_item = []
    for key in  c1.keys():
        for value in extracted_file1:
            if value[0] == key:
                #print(type(value[2]))
                extracted_file2.append(value)
                count_item.append(value)
                continue
    #print(count_item)
    #print()

    print(len(extracted_file2))
    data2 = pd.DataFrame(extracted_file2)
    data2.to_excel('newlist8.xlsx')
    #for no in extracted_file2:
        #print(no[2])
    import collections
    import matplotlib.pyplot as plt
    l = ['a', 'b', 'b', 'b', 'c']
    # w = collections.Counter(l)
    #w = c
    #plt.bar(w.keys(), w.values())
    #plt.show()



counting()



def ppt1():
    from pptx import Presentation

    X = Presentation()
    Layout = X.slide_layouts[0]
    first_slide = X.slides.add_slide(Layout)  # Adding first slide
    first_slide.shapes.title.text = "Creating a powerpoint using Python"
    first_slide.placeholders[1].text = "Created by Tutorialpoints"
    from pptx import Presentation
    from pptx.util import Inches
    Second_Layout = X.slide_layouts[5]
    second_slide = X.slides.add_slide(Second_Layout)
    second_slide.shapes.title.text = "Second slide"
    textbox = second_slide.shapes.add_textbox(Inches(3), Inches(1.5), Inches(3), Inches(1))
    textframe = textbox.text_frame
    paragraph = textframe.add_paragraph()
    paragraph.text = "This is a paragraph in the second slide!"

    X.save("First_presentation4.pptx")
#ppt1()

def ppt2():
    from pptx import Presentation
    from pptx.util import Inches
    X = Presentation()
    Layout = X.slide_layouts[0]
    first_slide = X.slides.add_slide(Layout)
    first_slide.shapes.title.text = "Creating a powerpoint using Python"
    first_slide.placeholders[1].text = "Created by Tutorialpoints"
    X.save("First_presentation.pptx")
    Second_Layout = X.slide_layouts[5]
    second_slide = X.slides.add_slide(Second_Layout)
    second_slide.shapes.title.text = "Second slide"
    textbox = second_slide.shapes.add_textbox(Inches(3), Inches(1.5), Inches(3), Inches(1))
    textframe = textbox.text_frame
    paragraph = textframe.add_paragraph()
    paragraph.text = "This is a paragraph in the second slide!"
    X.save("First_presentation.pptx")
#ppt2()


